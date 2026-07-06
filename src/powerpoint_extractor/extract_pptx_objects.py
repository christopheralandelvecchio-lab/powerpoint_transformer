"""
extract_pptx_objects.py

Purpose:
    Reusable object extraction library for reading native PowerPoint (.pptx)
    objects and converting them into structured SlideObject records.

Scope:
    Phase 1 object detection and extraction only. This module does not infer
    relationships between objects.

Primary Interface:
    extract_pptx_objects(pptx_path: Path) -> list[SlideObject]

Design Note:
    Command-line interaction is handled by main.py. Relationship extraction is
    handled by extract_pptx_relationships.py.
"""

from __future__ import annotations

from pathlib import Path
from typing import Iterable, Optional

from pptx import Presentation

from .models import SlideObject
from .utils import emu_to_inches, safe_text, shape_type_name


def extract_shapes(
    shapes: Iterable,
    slide_number: int,
    parent_group_id: Optional[str] = None,
    object_counter: Optional[list[int]] = None,
) -> list[SlideObject]:
    """
    Recursively extract objects from a slide or group of shapes.

    Parameters
    ----------
    shapes:
        Collection of PowerPoint shapes to inspect.

    slide_number:
        One-based slide number currently being processed.

    parent_group_id:
        Object ID of the parent group when these shapes are inside a group.
        None means the shapes are directly on the slide.

    object_counter:
        Shared mutable counter used to create unique object IDs. A list is used
        so recursive calls can update the same counter.

    Returns
    -------
    list[SlideObject]
        Flat list of extracted PowerPoint object records.

    Key algorithm points
    --------------------
    1. Loop through each shape in the current shape collection.
    2. Assign a unique object ID.
    3. Extract object metadata and geometry.
    4. Add the object to the result list.
    5. If the shape is a group, recursively extract its child shapes.
    """

    if object_counter is None:
        object_counter = [0]

    objects: list[SlideObject] = []

    for shape in shapes:
        object_counter[0] += 1
        object_id = f"S{slide_number:03d}_OBJ_{object_counter[0]:04d}"

        shape_type = shape_type_name(shape)

        obj = SlideObject(
            object_id=object_id,
            slide_number=slide_number,
            parent_group_id=parent_group_id,
            shape_name=getattr(shape, "name", ""),
            shape_type=shape_type,
            text=safe_text(shape),
            left_in=emu_to_inches(getattr(shape, "left", None)),
            top_in=emu_to_inches(getattr(shape, "top", None)),
            width_in=emu_to_inches(getattr(shape, "width", None)),
            height_in=emu_to_inches(getattr(shape, "height", None)),
            rotation=getattr(shape, "rotation", None),
        )

        objects.append(obj)

        # Grouped shapes expose a nested shape collection. Recursion preserves
        # containment through parent_group_id while still producing a flat list
        # of SlideObject records for CSV export and downstream ingestion.
        if shape_type == "GROUP":
            objects.extend(
                extract_shapes(
                    shape.shapes,
                    slide_number=slide_number,
                    parent_group_id=object_id,
                    object_counter=object_counter,
                )
            )

    return objects


def extract_pptx_objects(pptx_path: Path) -> list[SlideObject]:
    """
    Extract all native PowerPoint objects from a presentation.

    Parameters
    ----------
    pptx_path:
        Path to the .pptx file to process.

    Returns
    -------
    list[SlideObject]
        List of all objects detected across all slides.

    Raises
    ------
    FileNotFoundError
        If the input file does not exist.

    ValueError
        If the input file is not a .pptx file.
    """

    pptx_path = Path(pptx_path)

    if not pptx_path.exists():
        raise FileNotFoundError(f"PowerPoint file not found: {pptx_path}")

    if pptx_path.suffix.lower() != ".pptx":
        raise ValueError("Input file must be a .pptx file.")

    presentation = Presentation(pptx_path)
    all_objects: list[SlideObject] = []

    for index, slide in enumerate(presentation.slides, start=1):
        all_objects.extend(
            extract_shapes(
                slide.shapes,
                slide_number=index,
            )
        )

    return all_objects
