"""
utils.py

Purpose:
    Shared utility functions used by the PowerPoint extraction package.

Scope:
    This module contains reusable normalization and geometry helpers. It does
    not own object extraction, relationship extraction, or CSV writing.
"""

from __future__ import annotations

from typing import Optional

from pptx.enum.shapes import MSO_SHAPE_TYPE


EMU_PER_INCH = 914400


def emu_to_inches(value) -> Optional[float]:
    """
    Convert a PowerPoint EMU value to inches.

    Parameters
    ----------
    value:
        Numeric PowerPoint measurement in English Metric Units.

    Returns
    -------
    Optional[float]
        Measurement converted to inches and rounded to three decimals.
        Returns None when the input value is None.
    """

    if value is None:
        return None

    return round(int(value) / EMU_PER_INCH, 3)


def safe_text(shape) -> str:
    """
    Extract text from a PowerPoint shape, if text exists.

    Parameters
    ----------
    shape:
        PowerPoint shape object from python-pptx.

    Returns
    -------
    str
        Text contained in the shape. Returns an empty string when the shape has
        no text frame.

    Key algorithm points
    --------------------
    1. Check whether the shape has a text frame.
    2. Walk through paragraphs.
    3. Walk through runs inside each paragraph.
    4. Join non-blank paragraph text using pipe separators.
    """

    if not getattr(shape, "has_text_frame", False):
        return ""

    parts: list[str] = []

    for paragraph in shape.text_frame.paragraphs:
        paragraph_text = "".join(run.text for run in paragraph.runs).strip()

        if paragraph_text:
            parts.append(paragraph_text)

    return " | ".join(parts)


def shape_type_name(shape) -> str:
    """
    Convert a PowerPoint shape type value into a readable name.

    Parameters
    ----------
    shape:
        PowerPoint shape object from python-pptx.

    Returns
    -------
    str
        Human-readable shape type name when available. Falls back to the raw
        shape type value if conversion fails.
    """

    try:
        return MSO_SHAPE_TYPE(shape.shape_type).name
    except Exception:
        return str(shape.shape_type)
